#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
文件內容提取工具 - 支援 PPTX 和 PDF 格式
包含 OCR 圖片文字辨識功能
"""

import argparse
import json
import os
import sys
import io
import re
from pathlib import Path


def extract_image_text_ocr(image_bytes):
    """使用 OCR 從圖片提取文字"""
    try:
        from PIL import Image
        import pytesseract
    except ImportError:
        return ""
    
    try:
        image = Image.open(io.BytesIO(image_bytes))
        # 使用繁體中文 + 英文進行 OCR
        text = pytesseract.image_to_string(image, lang='chi_tra+eng')
        return text.strip()
    except Exception as e:
        return ""


def extract_table_text(shape):
    """從表格 shape 提取文字"""
    try:
        if not shape.has_table:
            return []
        
        table = shape.table
        rows_text = []
        for row in table.rows:
            row_cells = []
            for cell in row.cells:
                cell_text = cell.text.strip() if cell.text else ""
                row_cells.append(cell_text)
            if any(row_cells):
                rows_text.append(" | ".join(row_cells))
        return rows_text
    except:
        return []


def extract_pptx(file_path: str, use_ocr: bool = True) -> str:
    """從 PPTX 檔案提取文字內容，包含圖片 OCR"""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        print("請先安裝 python-pptx: pip install python-pptx")
        sys.exit(1)
    
    prs = Presentation(file_path)
    content = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_content = []
        image_texts = []
        
        for shape in slide.shapes:
            # 提取表格
            if shape.has_table:
                table_rows = extract_table_text(shape)
                for row in table_rows:
                    slide_content.append(row)
            
            # 提取文字方塊
            elif hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text.strip())
            
            # 提取圖片並進行 OCR
            if use_ocr:
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = shape.image
                        image_bytes = image.blob
                        ocr_text = extract_image_text_ocr(image_bytes)
                        if ocr_text and len(ocr_text) > 5:  # 過濾太短的 OCR 結果
                            image_texts.append(ocr_text)
                except:
                    pass
        
        if slide_content or image_texts:
            content.append(f"### 投影片 {slide_num}")
            
            # 先加入文字內容
            for text in slide_content:
                lines = text.split('\n')
                for line in lines:
                    line = line.strip()
                    if line:
                        content.append(f"- {line}")
            
            # 再加入 OCR 內容
            if image_texts:
                content.append("")
                content.append("**[圖片內容]**")
                for ocr_text in image_texts:
                    lines = ocr_text.split('\n')
                    for line in lines:
                        line = line.strip()
                        if line and len(line) > 1:
                            content.append(f"- {line}")
            
            content.append("")
    
    return "\n".join(content)


def extract_pdf(file_path: str, use_ocr: bool = True) -> str:
    """從 PDF 檔案提取文字內容，包含圖片 OCR"""
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        print("請先安裝 PyPDF2: pip install pypdf2")
        sys.exit(1)
    
    reader = PdfReader(file_path)
    content = []
    
    for page_num, page in enumerate(reader.pages, 1):
        page_content = []
        
        # 提取文字
        text = page.extract_text()
        if text and text.strip():
            lines = text.strip().split('\n')
            for line in lines:
                line = line.strip()
                if line:
                    page_content.append(f"- {line}")
        
        # 嘗試提取圖片並 OCR
        if use_ocr:
            try:
                if '/XObject' in page['/Resources']:
                    xObject = page['/Resources']['/XObject'].get_object()
                    for obj in xObject:
                        if xObject[obj]['/Subtype'] == '/Image':
                            try:
                                data = xObject[obj].get_data()
                                ocr_text = extract_image_text_ocr(data)
                                if ocr_text and len(ocr_text) > 5:
                                    page_content.append("")
                                    page_content.append("**[圖片內容]**")
                                    for line in ocr_text.split('\n'):
                                        line = line.strip()
                                        if line and len(line) > 1:
                                            page_content.append(f"- {line}")
                            except:
                                pass
            except:
                pass
        
        if page_content:
            content.append(f"### 第 {page_num} 頁")
            content.extend(page_content)
            content.append("")
    
    return "\n".join(content)


def extract_document(file_path: str, use_ocr: bool = True) -> tuple:
    """
    根據檔案類型提取內容
    返回: (檔案名稱, 內容)
    """
    path = Path(file_path)
    file_name = path.stem
    
    # 移除 UUID 後綴 (如果有的話)
    uuid_pattern = r'-[a-f0-9]{8}-[a-f0-9]{4}-[a-f0-9]{4}-[a-f0-9]{4}-[a-f0-9]{12}$'
    file_name = re.sub(uuid_pattern, '', file_name)
    
    ext = path.suffix.lower()
    
    if ext == '.pptx':
        return file_name, extract_pptx(file_path, use_ocr)
    elif ext == '.pdf':
        return file_name, extract_pdf(file_path, use_ocr)
    else:
        return file_name, f"不支援的檔案格式: {ext}"


def find_department(filename: str, dept_mapping: dict) -> str:
    """根據檔案名稱找出對應的部門"""
    for dept, keywords in dept_mapping.items():
        for keyword in keywords:
            if keyword in filename:
                return dept
    return "其他"


def process_folder(folder_path: str, dept_mapping: dict = None, use_ocr: bool = True) -> str:
    """處理資料夾中的所有文件"""
    folder = Path(folder_path)
    files = list(folder.glob("*.pptx")) + list(folder.glob("*.pdf"))
    
    if not files:
        return "資料夾中沒有找到 PPTX 或 PDF 檔案"
    
    print(f"找到 {len(files)} 個檔案，開始處理...")
    
    # 如果有部門對應，則按部門分組
    if dept_mapping:
        dept_contents = {}
        for i, file in enumerate(files, 1):
            print(f"處理中 [{i}/{len(files)}]: {file.name}")
            file_name, content = extract_document(str(file), use_ocr)
            dept = find_department(file_name, dept_mapping)
            
            if dept not in dept_contents:
                dept_contents[dept] = []
            
            dept_contents[dept].append({
                "name": file_name,
                "content": content
            })
        
        # 組合輸出
        output = []
        for dept in dept_mapping.keys():
            if dept in dept_contents:
                output.append(f"# {dept}\n")
                for doc in dept_contents[dept]:
                    output.append(f"## {doc['name']}\n")
                    output.append(doc['content'])
                    output.append("\n---\n")
        
        # 處理未分類的文件
        if "其他" in dept_contents:
            output.append("# 其他\n")
            for doc in dept_contents["其他"]:
                output.append(f"## {doc['name']}\n")
                output.append(doc['content'])
                output.append("\n---\n")
        
        return "\n".join(output)
    else:
        output = []
        for i, file in enumerate(files, 1):
            print(f"處理中 [{i}/{len(files)}]: {file.name}")
            file_name, content = extract_document(str(file), use_ocr)
            output.append(f"## {file_name}\n")
            output.append(content)
            output.append("\n---\n")
        return "\n".join(output)


def main():
    parser = argparse.ArgumentParser(
        description='從 PPTX/PDF 檔案提取文字內容 (支援 OCR)'
    )
    parser.add_argument('path', help='檔案或資料夾路徑')
    parser.add_argument('-o', '--output', help='輸出檔案路徑')
    parser.add_argument('--group-by-dept', help='部門對應 JSON 檔案路徑')
    parser.add_argument('--no-ocr', action='store_true', help='停用圖片 OCR 功能')
    
    args = parser.parse_args()
    
    path = Path(args.path)
    use_ocr = not args.no_ocr
    
    if use_ocr:
        # 檢查 OCR 依賴
        try:
            import pytesseract
            from PIL import Image
            print("OCR 功能已啟用 (使用 Tesseract)")
        except ImportError:
            print("警告: OCR 依賴未安裝，將只提取純文字")
            print("安裝 OCR: pip install pytesseract pillow")
            use_ocr = False
    
    if path.is_file():
        file_name, content = extract_document(str(path), use_ocr)
        result = f"## {file_name}\n\n{content}"
    elif path.is_dir():
        dept_mapping = None
        if args.group_by_dept:
            with open(args.group_by_dept, 'r', encoding='utf-8') as f:
                dept_mapping = json.load(f)
        result = process_folder(str(path), dept_mapping, use_ocr)
    else:
        print(f"找不到路徑: {path}")
        sys.exit(1)
    
    if args.output:
        with open(args.output, 'w', encoding='utf-8') as f:
            f.write(result)
        print(f"已輸出到: {args.output}")
    else:
        print(result)


if __name__ == '__main__':
    main()
