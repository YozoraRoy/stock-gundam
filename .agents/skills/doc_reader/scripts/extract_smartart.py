#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
進階 PPTX 內容提取 - 包含 SmartArt 和圖表
直接解析 PPTX 內的 XML 結構
"""

import argparse
import json
import os
import sys
import re
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET


# XML 命名空間
NAMESPACES = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'dgm': 'http://schemas.openxmlformats.org/drawingml/2006/diagram',
    'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart',
}


def extract_text_from_xml(element) -> list:
    """遞迴提取 XML 元素中的所有文字"""
    texts = []
    
    # 取得當前元素的文字
    if element.text and element.text.strip():
        texts.append(element.text.strip())
    
    # 遞迴處理子元素
    for child in element:
        texts.extend(extract_text_from_xml(child))
    
    # 取得 tail 文字
    if element.tail and element.tail.strip():
        texts.append(element.tail.strip())
    
    return texts


def extract_smartart_from_pptx(pptx_path: str) -> dict:
    """
    從 PPTX 中提取 SmartArt 內容
    返回: {slide_num: [texts]}
    """
    smartart_content = {}
    
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            # 列出所有檔案
            file_list = zf.namelist()
            
            # 找出 SmartArt 資料檔案
            diagram_files = [f for f in file_list if 'diagrams/data' in f.lower()]
            
            for diagram_file in diagram_files:
                try:
                    with zf.open(diagram_file) as f:
                        tree = ET.parse(f)
                        root = tree.getroot()
                        texts = extract_text_from_xml(root)
                        
                        # 過濾空白和太短的文字
                        texts = [t for t in texts if len(t) > 1]
                        
                        if texts:
                            # 從檔名推測投影片編號
                            match = re.search(r'data(\d+)', diagram_file)
                            slide_num = int(match.group(1)) if match else 0
                            
                            if slide_num not in smartart_content:
                                smartart_content[slide_num] = []
                            smartart_content[slide_num].extend(texts)
                except Exception as e:
                    print(f"  處理 {diagram_file} 時發生錯誤: {e}")
            
            # 也檢查圖表資料
            chart_files = [f for f in file_list if 'charts/chart' in f.lower()]
            
            for chart_file in chart_files:
                try:
                    with zf.open(chart_file) as f:
                        tree = ET.parse(f)
                        root = tree.getroot()
                        texts = extract_text_from_xml(root)
                        texts = [t for t in texts if len(t) > 1]
                        
                        if texts:
                            match = re.search(r'chart(\d+)', chart_file)
                            chart_num = int(match.group(1)) if match else 0
                            
                            key = f"chart_{chart_num}"
                            smartart_content[key] = texts
                except Exception as e:
                    print(f"  處理 {chart_file} 時發生錯誤: {e}")
    
    except Exception as e:
        print(f"無法讀取 PPTX: {e}")
    
    return smartart_content


def extract_pptx_full(file_path: str) -> str:
    """
    完整提取 PPTX 內容，包含 SmartArt
    """
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        print("請先安裝 python-pptx: pip install python-pptx")
        sys.exit(1)
    
    path = Path(file_path)
    file_name = path.stem
    
    # 移除 UUID 後綴
    uuid_pattern = r'-[a-f0-9]{8}-[a-f0-9]{4}-[a-f0-9]{4}-[a-f0-9]{4}-[a-f0-9]{12}$'
    file_name = re.sub(uuid_pattern, '', file_name)
    
    # 提取 SmartArt 內容
    smartart_content = extract_smartart_from_pptx(file_path)
    
    # 使用 python-pptx 提取一般文字
    prs = Presentation(file_path)
    content = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_content = []
        
        # 提取表格
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_cells = []
                    for cell in row.cells:
                        if cell.text and cell.text.strip():
                            row_cells.append(cell.text.strip())
                    if row_cells:
                        slide_content.append(" | ".join(row_cells))
            
            # 提取文字方塊
            elif hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text.strip())
        
        # 加入 SmartArt 內容
        smartart_texts = smartart_content.get(slide_num, [])
        
        if slide_content or smartart_texts:
            content.append(f"### 投影片 {slide_num}")
            
            # 一般內容
            for text in slide_content:
                lines = text.split('\n')
                for line in lines:
                    line = line.strip()
                    if line:
                        content.append(f"- {line}")
            
            # SmartArt 內容
            if smartart_texts:
                content.append("")
                content.append("**[SmartArt/圖表內容]**")
                for text in smartart_texts:
                    content.append(f"- {text}")
            
            content.append("")
    
    # 加入未對應的圖表內容
    for key, texts in smartart_content.items():
        if isinstance(key, str) and key.startswith("chart_"):
            content.append(f"### 圖表資料")
            for text in texts:
                content.append(f"- {text}")
            content.append("")
    
    return "\n".join(content)


def find_department(filename: str, dept_mapping: dict) -> str:
    """根據檔案名稱找出對應的部門"""
    for dept, keywords in dept_mapping.items():
        for keyword in keywords:
            if keyword in filename:
                return dept
    return "其他"


def process_file(file_path: str) -> tuple:
    """處理單一檔案"""
    path = Path(file_path)
    file_name = path.stem
    
    uuid_pattern = r'-[a-f0-9]{8}-[a-f0-9]{4}-[a-f0-9]{4}-[a-f0-9]{4}-[a-f0-9]{12}$'
    file_name = re.sub(uuid_pattern, '', file_name)
    
    ext = path.suffix.lower()
    
    if ext == '.pptx':
        return file_name, extract_pptx_full(str(path))
    elif ext == '.pdf':
        # 對 PDF 使用基本提取
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            content = []
            for page_num, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    content.append(f"### 第 {page_num} 頁")
                    for line in text.strip().split('\n'):
                        line = line.strip()
                        if line:
                            content.append(f"- {line}")
                    content.append("")
            return file_name, "\n".join(content)
        except ImportError:
            return file_name, "需要安裝 PyPDF2"
    else:
        return file_name, f"不支援的格式: {ext}"


def process_folder(folder_path: str, dept_mapping: dict = None) -> str:
    """處理資料夾"""
    folder = Path(folder_path)
    files = list(folder.glob("*.pptx")) + list(folder.glob("*.pdf"))
    
    if not files:
        return "資料夾中沒有找到 PPTX 或 PDF 檔案"
    
    print(f"找到 {len(files)} 個檔案，開始處理...")
    
    if dept_mapping:
        dept_contents = {}
        for i, file in enumerate(files, 1):
            print(f"處理中 [{i}/{len(files)}]: {file.name}")
            file_name, content = process_file(str(file))
            dept = find_department(file_name, dept_mapping)
            
            if dept not in dept_contents:
                dept_contents[dept] = []
            
            dept_contents[dept].append({
                "name": file_name,
                "content": content
            })
        
        output = []
        for dept in dept_mapping.keys():
            if dept in dept_contents:
                output.append(f"# {dept}\n")
                for doc in dept_contents[dept]:
                    output.append(f"## {doc['name']}\n")
                    output.append(doc['content'])
                    output.append("\n---\n")
        
        if "其他" in dept_contents:
            output.append("# 其他\n")
            for doc in dept_contents["其他"]:
                output.append(f"## {doc['name']}\n")
                output.append(doc['content'])
                output.append("\n---\n")
        
        return "\n".join(output)
    else:
        output = []
        for i, file in enumerate(files, 1):
            print(f"處理中 [{i}/{len(files)}]: {file.name}")
            file_name, content = process_file(str(file))
            output.append(f"## {file_name}\n")
            output.append(content)
            output.append("\n---\n")
        return "\n".join(output)


def main():
    parser = argparse.ArgumentParser(
        description='完整 PPTX 內容提取 (含 SmartArt)'
    )
    parser.add_argument('path', help='檔案或資料夾路徑')
    parser.add_argument('-o', '--output', help='輸出檔案路徑')
    parser.add_argument('--group-by-dept', help='部門對應 JSON 檔案路徑')
    
    args = parser.parse_args()
    path = Path(args.path)
    
    if path.is_file():
        file_name, content = process_file(str(path))
        result = f"## {file_name}\n\n{content}"
    elif path.is_dir():
        dept_mapping = None
        if args.group_by_dept:
            with open(args.group_by_dept, 'r', encoding='utf-8') as f:
                dept_mapping = json.load(f)
        result = process_folder(str(path), dept_mapping)
    else:
        print(f"找不到路徑: {path}")
        sys.exit(1)
    
    if args.output:
        with open(args.output, 'w', encoding='utf-8') as f:
            f.write(result)
        print(f"已輸出到: {args.output}")
    else:
        print(result)


if __name__ == '__main__':
    main()
